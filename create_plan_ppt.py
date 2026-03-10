from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_presentation():
    prs = Presentation()
    
    # --- Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "プロジェクト計画案（要件管理）"
    subtitle.text = "要件ブレを防ぎ、計画的に進めるためのアプローチ\n2026年3月"

    # --- Slide 1: 概要 ---
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    title_shape.text = "要件管理の全体像"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "新規プロジェクトにおける要件の不確実性をコントロールする4つのステップ\n"
    
    p = tf.add_paragraph()
    p.text = "1. スコープの境界線を明確にする（定義フェーズ）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. 不確実性の高い部分を早期に検証する（設計検証フェーズ）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. 変更管理プロセスを確立する（運用ルール）"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. 短いサイクルで確認とフィードバックを行う（開発フェーズ）"
    p.level = 1

    # --- Slide 2: ステップ1 ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "1. スコープの境界線を明確にする"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "要件がブレる最大の原因は、関係者間で「どこまでやるか」の認識が揃っていないこと"
    
    p = tf.add_paragraph()
    p.text = "プロジェクトの目的とゴール"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "何の課題を解決するためのシステムかを簡潔に定義（判断の原点）"
    p2.level = 2
    
    p = tf.add_paragraph()
    p.text = "やらないこと（アウト・オブ・スコープ）の定義"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "「今回はやらないこと」を明記し、要件の肥大化を防ぐ防波堤とする"
    p2.level = 2

    p = tf.add_paragraph()
    p.text = "要求の優先順位付け"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Must（必須）、Should（推奨）、Could（できれば）、Wont（今回はやらない）の分類"
    p2.level = 2

    # --- Slide 3: ステップ2 ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "2. 不確実性の高い部分を早期に検証する"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "言葉だけの議論による認識のズレを防ぐため、早期に視覚化・検証を行う"
    
    p = tf.add_paragraph()
    p.text = "ワイヤーフレームやプロトタイプの作成"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "UIや主要な業務フローの画面レイアウト・遷移図を作成し、ユーザー目線でレビュー"
    p2.level = 2
    
    p = tf.add_paragraph()
    p.text = "不明瞭な要件の洗い出しとPoC"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "技術的実現性が不安な部分や仕様未定部分を特定し、小規模な実装検証を先行"
    p2.level = 2

    # --- Slide 4: ステップ3 ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "3. 変更管理プロセスを確立する"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "開発途中の要望追加・変更を無秩序に受け入れないための運用ルール"
    
    p = tf.add_paragraph()
    p.text = "変更要求の受付フロー化"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "口頭ベースで受けず、必ず課題管理ツール（Issue等）に起票して可視化"
    p2.level = 2
    
    p = tf.add_paragraph()
    p.text = "トレードオフの事前合意"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "新要件追加時は「スケジュール延長」「別要件削減」「コスト追加」のいずれかが必要との原則を事前合意"
    p2.level = 2

    # --- Slide 5: ステップ4 ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "4. 短いサイクルで確認とフィードバックを行う"
    
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "長期間作り込んだ後の「思っていたものと違う」を防ぐ"
    
    p = tf.add_paragraph()
    p.text = "定期的なデモンストレーション"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "1〜2週間の短いサイクル（スプリント）で、完成した「動く画面・機能」を定期確認"
    p2.level = 2
    
    p = tf.add_paragraph()
    p.text = "軌道修正の定常化"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "デモを通じて得たフィードバックを次のサイクルに組み込み、小さな方向修正を繰り返す"
    p2.level = 2
    
    # Save presentation
    output_path = os.path.join("c:/Users/owner/Desktop/一時/Google_Antigravity_Proj", "プロジェクト計画案_要件管理.pptx")
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == '__main__':
    create_presentation()
